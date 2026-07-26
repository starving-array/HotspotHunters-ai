from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation(r'E:\project\HotspotHunters-ai\KSP Datathon 2026 _ Prototype Submission Template.pptx')

print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")
print(f"Total slides: {len(prs.slides)}")
print()

for i, slide in enumerate(prs.slides):
    print(f"{'='*60}")
    print(f"SLIDE {i+1}: Layout=\"{slide.slide_layout.name}\"")
    print(f"{'='*60}")
    for shape in slide.shapes:
        txt = shape.text[:300] if shape.has_text_frame else "N/A"
        print(f"  [{shape.shape_type}] Name=\"{shape.name}\"")
        print(f"    Left={shape.left}, Top={shape.top}, Width={shape.width}, Height={shape.height}")
        if shape.has_table:
            table = shape.table
            for r_idx, row in enumerate(table.rows):
                cells = [cell.text for cell in row.cells]
                print(f"    Table Row {r_idx}: {cells}")
        print(f"    Text: \"{txt}\"")
        print()
    print()

# Also check slide layouts available
print("\n\nSLIDE LAYOUTS:")
for i, layout in enumerate(prs.slide_layouts):
    print(f"  Layout {i}: \"{layout.name}\"")
    for ph in layout.placeholders:
        print(f"    Placeholder {ph.placeholder_format.idx}: {ph.name} ({ph.placeholder_format.type})")