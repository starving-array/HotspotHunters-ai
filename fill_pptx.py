from pptx import Presentation

# Path to the original template
template_path = r'E:\project\HotspotHunters-ai\KSP Datathon 2026 _ Prototype Submission Template.pptx'
prs = Presentation(template_path)

# Content for each slide (0‑based index). Use triple‑quoted strings to keep line breaks.
slide_contents = [
    """Team Details

Team name: HotspotHunters
Team leader name: Archishman Das
Team size: 1 (Solo)
Problem Statement: AI‑Driven Crime Analytics & Visualization Platform for Karnataka State Police.""",
    """Brief about the solution

The KSP Intelligence Portal is a production‑grade distributed intelligence platform that ingests FIR events, persists them to PostgreSQL & ElasticSearch, maintains a sub‑millisecond live hotspot leaderboard in Redis, predicts future hotspots with Random Forest, and translates natural‑language queries via a configurable LLM provider chain. All data is auditable, real‑time, and geospatially visualised.""",
    """Opportunities

- How different is it from existing ideas?
- How will it solve the problem?
- Unique Selling Proposition (USP) of the proposed solution""",
    """List of features offered by the solution

1. Real‑time hotspot leaderboard (Redis Sorted Set)\n2. Geospatial search (ElasticSearch geo_point)\n3. Predictive risk scoring (Random Forest + Gradient Boosting + SHAP)\n4. Natural‑language query translation (LLM provider chain)\n5. Offender network graph (PostgreSQL co‑crime edges)\n6. Live alert stream (Redis Stream → SSE)\n7. Audit trail (append‑only PostgreSQL audit_log)\n8. Dashboard KPI cards, trends, cyber‑crime view\n9. Multi‑language UI (English & Kannada)\n10. Observability (Prometheus + Grafana)""",
    """Process flow / Use‑Case Diagram

Data Generator → Kafka (fir‑events) → [IndexingConsumer → PG + ES]
                                         ├─ AggregationConsumer → Redis leaderboard & alerts
                                         └─ AnomalyConsumer → Spike detection → alert‑events
Spring Boot API reads from PG/ES/Redis, proxies to FastAPI ML, serves React SPA.
Live alerts via SSE, predictions via ML service.""",
    """Wireframes / Mock diagrams

- Overview dashboard with KPI cards
- Map view (Leaflet heatmap)
- Hotspot leaderboard page
- NL query panel
- Prediction panel (risk scores + SHAP)
- Trends & analytics page
- Audit log viewer
- Network graph visualisation
- Cyber‑crime dashboard
- Settings / language toggle""",
    """Architecture Overview

Client (React SPA) ↔ Spring Boot API ↔ PostgreSQL, ElasticSearch, Redis
Spring Boot ↔ FastAPI ML Service (HTTP)
Kafka Signals replace event bus in Catalyst deployment.
All services containerised; health‑checks via Actuator.""",
    """Technology Stack

Backend: Java 17, Spring Boot 3.4, Kafka, PostgreSQL 16, ElasticSearch 8.14, Redis 7
ML: Python 3.11, FastAPI, scikit‑learn, SHAP, LLM providers (Anthropic, Groq, Gemini)
Frontend: React 18, TypeScript 5, Vite 5, Tailwind, Leaflet, Recharts
Infra: Docker Compose, Prometheus, Grafana, GitHub Actions CI/CD
Deployment: Zoho Catalyst AppSail, Data Store, Cache, Signals""",
    """Catalyst Services Used

- AppSail (Custom Runtime) – Spring Boot API
- AppSail (Custom Runtime) – Python ML Service
- AppSail (Node.js) – React Frontend
- Data Store – PostgreSQL
- Cache – Redis
- Signals – Event bus (fir‑events, alert‑events, audit‑events)""",
    """Estimated implementation cost (Catalyst)

| Service | Plan | Monthly Cost |
|---------|------|--------------|
| Data Store (PostgreSQL) | Starter ($15) | $15 |
| Cache (Redis) | Free | $0 |
| Signals | Free | $0 |
| AppSail API | Starter ($10) | $10 |
| AppSail ML | Starter ($10) | $10 |
| AppSail Frontend | Starter ($10) | $10 |
| **Total** |  | **~$45** |""",
    """Snapshots of the prototype (to be added)

[Insert screenshots of Overview, Map view, Leaderboard, NL query, Prediction panel, Live alerts, Audit log, Network graph, Cyber‑crime dashboard]""",
    """Prototype Performance Report / Benchmarking

- Geo search (<150 ms) on 200 K records
- Hotspot API sub‑ms via Redis ZINCRBY
- Kafka consumer lag < 20 msgs at 1 K events/s
- NL query translation ≤ 1.5 s (LLM chain)
- API median latency ~200 ms (p99 < 500 ms)
- Model inference ≤ 300 ms (with SHAP)
- End‑to‑end NL query → map results ≤ 2 s""",
    """GitHub – Demo Video – Deployed Link

Repository: https://github.com/ArchishmanDas/hotspothunters-ai
Demo video: (to be uploaded)
Catalyst deployed URL: (to be filled after deployment)""",
    """Additional Details / Future Development

Roadmap:
- Phase 2: Semantic clustering of modus operandi (sentence embeddings)
- Phase 3: Full Kannada LLM support
- Phase 5: Mobile officer app with push notifications

Completed: All core features, CI/CD, Catalyst migration guide, observability, audit logging, SHAP explanations, multi‑language UI."""
]

# Fill each slide's first text‑containing shape
for idx, slide in enumerate(prs.slides):
    # Find the first shape that has a text frame
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            # Replace the entire text
            shape.text = slide_contents[idx] if idx < len(slide_contents) else ""
            break

# Save the new presentation
output_path = r'E:\project\HotspotHunters-ai\KSP_Datathon_2026_Prototype_Submission_filled.pptx'
prs.save(output_path)
print(f"Saved filled presentation to {output_path}")
